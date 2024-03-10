import openai
from openai import OpenAI
import os
import threading
import time
from datetime import datetime
from tqdm import tqdm
from pptx import Presentation
from pptx.exc import PackageNotFoundError

client = OpenAI() # Create an instance of the OpenAI class

# Set up the OpenAI API
openai.api_key = os.getenv("OPENAI_API_KEY")

current_datetime = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
presentation_name = f"Cybersecurity_Awareness_Training_{current_datetime}"

def content_to_slide(slide_content: str, prs):
    try:
        slide_layout = prs.slide_layouts[1]  # Use a slide layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        title, content_and_notes = slide_content.split('\n', 1)
        content, speaker_notes = content_and_notes.split('---', 1)
        
        slide.shapes.title.text = title.strip()  # Setting the title

        content_placeholder = slide.placeholders[1]  # Adding the slide content
        for line in content.strip().split('\n'):
            p = content_placeholder.text_frame.add_paragraph()
            p.text = line.strip()
            
        notes_slide = slide.notes_slide  # Adding the speaker notes
        text_frame = notes_slide.notes_text_frame
        text_frame.text = speaker_notes.strip()

    except Exception as e:
        print(f"An error occurred while creating the PowerPoint slide: {e}")
        return False
    return True

# Function to display elapsed time while waiting for the API call
def display_elapsed_time(event):
    start_time = time.time()
    while not event.is_set():
        elapsed_time = time.time() - start_time
        print(f"\rElapsed time: {elapsed_time:.2f} seconds", end="")
        time.sleep(1)

# Create an Event object
api_call_completed = threading.Event()

# Starting the thread for displaying elapsed time
elapsed_time_thread = threading.Thread(target=display_elapsed_time, args=(api_call_completed,))
elapsed_time_thread.start()

# Prepare initial prompt
messages=[
    {
        "role": "system",
        "content": "You are a cybersecurity professional with more than 25 years of experience."
    },
    {
        "role": "user",
        "content": "Create a cybersecurity awareness training slide list that will be used for a PowerPoint slide based awareness training course, for company employees, for the electric utility industry. This should be a single level list and should not contain subsections. Each item should represent a single slide."
    }
]

print(f"\nGenerating training outline...")
try:
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=2048,
        n=1,
        stop=None,
        temperature=0.7,
    )
except Exception as e:
    print("An error occurred while connecting to the OpenAI API:", e)
    exit(1)

# Get outline
outline = response.choices[0].message.content.strip()

print(outline + "\n")

# Split outline into sections
sections = outline.split("\n")

# The path to the template file
ppt_template_path = 'ppt_template.pptx'  # Replace with your actual path

# Use an existing PowerPoint presentation template or create a new one
try:
    prs = Presentation(ppt_template_path)
except Exception:
    prs = Presentation()

# For each section in the outline
for i, section in tqdm(enumerate(sections, start=1), total=len(sections), leave=False):
    print(f"\nGenerating details for section {i}...")

    # Prepare prompt for detailed info
    messages=[
        {
            "role": "system",
            "content": "You are a cybersecurity professional with more than 25 years of experience."
        },
        {
            "role": "user",
            "content": f"You are currently working on a PowerPoint presentation that will be used for a cybersecurity awareness training course, for end users, for the electric utility industry. The following outline is being used:\n\n{outline}\n\nCreate a single slide for the following section (and only this section) of the outline: {section}. The slides are for the employee's viewing, not the instructor, so use the appropriate voice and perspective. The employee will be using these slides as the primary source of information and lecture for the course. So, include the necessary lecture script in the speaker notes section. Do not write anything that should go in another section of the policy. Do not use the word 'Content' at the beginning of the content section. Do not use bullet characters or markdown. Use a tab for sub levels. The following format must be strictly followed:\n\n[Title]\n[Content]\n---\n[Lecture]"
        }
    ]

    # Reset the Event before each API call
    api_call_completed.clear()

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=messages,
            max_tokens=2048,
            n=1,
            stop=None,
            temperature=0.7,
        )
    except Exception as e:
        print("An error occurred while connecting to the OpenAI API:", e)
        exit(1)

    # Set the Event to signal that the API call is complete
    api_call_completed.set()

    # Get detailed info
    slide_content = response.choices[0].message.content.strip()

    # Create a slide with the detailed info
    if not content_to_slide(slide_content, prs):
        print("Failed to create slide. Skipping to next section...")
        continue

# Save the PowerPoint presentation
pptx_output_file = f"{presentation_name}_presentation.pptx"
try:
    prs.save(pptx_output_file)
    print("\nPresentation generated successfully!")
except PackageNotFoundError:
    print("\nAn error occurred while saving the PowerPoint presentation. Please check the file path and try again.")
except Exception as e:
    print(f"\nAn error occurred while saving the PowerPoint presentation: {e}")

# At the end of the script, make sure to join the elapsed_time_thread
api_call_completed.set()
elapsed_time_thread.join()
